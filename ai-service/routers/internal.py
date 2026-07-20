"""AI服务 - 文件解析路由模块

提供内部REST接口供Java后端调用，处理上传文件的文本提取。
大文件解析通过 run_in_executor 放入线程池，不阻塞事件循环。
"""

import asyncio
import base64
import logging
import os as _os
from typing import Dict

from fastapi import APIRouter, HTTPException
from pydantic import BaseModel

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/internal", tags=["file-parser"])


class ParseFileRequest(BaseModel):
    """文件解析请求"""
    file_path: str
    file_name: str


@router.get("/health")
async def health() -> Dict[str, str]:
    return {"status": "healthy", "service": "file-parser"}


@router.post("/parse-file")
async def parse_file(request: ParseFileRequest):
    """解析上传的文件为纯文本 — 支持 txt/md/docx/pdf/pptx/xlsx/csv/html/png/jpg/mp3/wav

    同步IO通过 run_in_executor 异步化，避免阻塞 uvicorn 事件循环。
    """
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _parse_sync, request.file_path, request.file_name)


def _parse_sync(file_path: str, file_name: str) -> dict:
    """同步文件解析逻辑，在线程池中运行。"""
    if not _os.path.exists(file_path):
        raise HTTPException(status_code=404, detail=f"文件不存在: {file_path}")

    ext = _os.path.splitext(file_name)[1].lower() if file_name else ""

    try:
        # 纯文本
        if ext in (".txt", ".md"):
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                text = f.read()
            return {"text": text, "file_type": ext, "size": len(text)}

        # Word
        elif ext == ".docx":
            from docx import Document
            doc = Document(file_path)
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            return {"text": text, "file_type": ext, "size": len(text)}

        # PDF
        elif ext == ".pdf":
            import PyPDF2
            reader = PyPDF2.PdfReader(file_path)
            text = "\n".join(p.extract_text() or "" for p in reader.pages)
            return {"text": text, "file_type": ext, "size": len(text)}

        # PowerPoint
        elif ext in (".ppt", ".pptx"):
            from pptx import Presentation
            prs = Presentation(file_path)
            slides = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        slides.append(shape.text_frame.text)
            text = "\n".join(slides)
            return {"text": text, "file_type": ext, "size": len(text)}

        # Excel
        elif ext in (".xls", ".xlsx", ".csv"):
            import pandas as pd
            if ext == ".csv":
                df = pd.read_csv(file_path)
            else:
                df = pd.read_excel(file_path)
            text = df.to_string(index=False)
            return {"text": text, "file_type": ext, "size": len(text)}

        # HTML
        elif ext in (".html", ".htm"):
            from bs4 import BeautifulSoup
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                soup = BeautifulSoup(f.read(), "html.parser")
            text = soup.get_text(separator="\n")
            return {"text": text, "file_type": ext, "size": len(text)}

        # 图片
        elif ext in (".png", ".jpg", ".jpeg"):
            with open(file_path, "rb") as f:
                img_b64 = base64.b64encode(f.read()).decode()
            return {
                "text": f"[图片文件: {file_name}]\n请使用多模态模型（如 GPT-4V / Gemini）解析此图片内容。\nBase64: data:image/{ext[1:]};base64,{img_b64[:200]}...",
                "file_type": ext,
                "size": _os.path.getsize(file_path),
                "needs_manual": True,
            }

        # 音频
        elif ext in (".mp3", ".wav"):
            with open(file_path, "rb") as f:
                audio_b64 = base64.b64encode(f.read()).decode()
            return {
                "text": f"[音频文件: {file_name}]\n请使用语音转文字服务（如 Whisper / 飞书妙记）转写此音频。\nBase64: data:audio/{ext[1:]};base64,{audio_b64[:200]}...",
                "file_type": ext,
                "size": _os.path.getsize(file_path),
                "needs_manual": True,
            }

        else:
            raise HTTPException(status_code=400, detail=f"不支持的文件格式: {ext}")

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"文件解析失败: {file_path} ({file_name}): {e}")
        raise HTTPException(status_code=500, detail=f"文件解析失败: {e}")
